import os
import shutil
from pptx import Presentation

def extract_prism_files(pptx_file, prism_folder):
    # Load the PowerPoint file
    presentation = Presentation(pptx_file)

    # Iterate over all slides in the presentation
    for slide_num, slide in enumerate(presentation.slides, start=1):
        # Iterate over all shapes in the slide
        for shape in slide.shapes:
            # Check if the shape is a Prism file
            if shape.shape_type == 14:  # 14 is the shape type for Embedded OLE Objects
                # Get the Prism file from the shape
                prism_file = shape.ole_format.embedded_package.data

                # Create a file path for the Prism file
                prism_file_path = os.path.join(prism_folder, f"Slide_{slide_num}.prism")

                # Write the Prism file to the Prism folder
                with open(prism_file_path, 'wb') as f:
                    f.write(prism_file)

from pptx import Presentation
import os

def extract_prism_files_from_pptx(pptx_path, output_folder='prism'):
    # Create the output folder if it doesn't exist
    os.makedirs(output_folder, exist_ok=True)

    # Load the PowerPoint presentation
    presentation = Presentation(pptx_path)

    # Iterate over all slides
    for slide_number, slide in enumerate(presentation.slides):
        # Iterate over all shapes in the slide
        for shape_number, shape in enumerate(slide.shapes):
            # Check if the shape has an OLE object (linked file)
            if shape.has_text_frame and 'prism' in shape.text_frame.text.lower():
                # Assuming the object contains the word 'prism' in the text
                prism_filename = f"{slide_number + 1}_{shape_number + 1}_prism.prism"

                # Save the Prism file to the output folder
                prism_path = os.path.join(output_folder, prism_filename)
                with open(prism_path, 'w') as prism_file:
                    # You might want to customize this part based on how the Prism file is stored in the PowerPoint
                    prism_file.write(shape.text_frame.text)

    print("Extraction complete.")


if __name__ == '__main__':
    # Create a folder to store the Prism files
    prism_folder = "Prism"
    if os.path.exists(prism_folder):
        shutil.rmtree(prism_folder)
    os.mkdir(prism_folder)

    # Extract the Prism files from the PowerPoint presentation
    extract_prism_files_from_pptx("Epilepsy-control figures slow.pptx", prism_folder)
