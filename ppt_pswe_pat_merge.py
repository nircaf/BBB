import pandas as pd
import os
from pathlib import Path

# import Presentation from pptx
from pptx import Presentation

# import mann whitney u test
from scipy.stats import mannwhitneyu

# matlibplot
import matplotlib.pyplot as plt
import base64
from PIL import Image
import sys
from io import StringIO

from bson.binary import Binary
from pymongo import MongoClient

# load data from excel and powerpoint files
Epilepsy_patients = pd.read_excel("Epilepsy_clinical_data.xlsx", sheet_name="All")
# drop row after index 28
Epilepsy_patients = Epilepsy_patients.drop(Epilepsy_patients.index[29:])
# epilepsy_age = Epilepsy_patients["'age'"]
# Controls = pd.read_excel('Epilepsy_clinical_data.xlsx', sheet_name='Controls')
# conrols_age = Controls["age'"]
# # do mann whitney u test on age
# print(mannwhitneyu(epilepsy_age, conrols_age))
# load pptx file
pptx_file = Presentation("Epilepsy_individulas.pptx")
pptx_dir = (
    r"D:\Dropbox (BGU DAL BBB Group)\Nir_Epilepsy_Controls\Python\Individual results"
)
# load Epilepsy all.csv
Epilepsy_all = pd.read_csv("Epilepsy all.csv")
Epilepsy_all["abs_diff"] = Epilepsy_all["abs_diff"].str.extract("(\d+)").astype(int)
# new column code2 which is code.split('_')[1]
Epilepsy_all["code2"] = Epilepsy_all["code"].apply(lambda x: x.split("_")[1])
df_eeg_closest = pd.DataFrame()
# run over epilepsy all groupby code2
for code, group in Epilepsy_all.groupby("code2"):
    min_abs_diff = group[group["abs_diff"] == group["abs_diff"].min()]
    # concat to df_eeg_closest the row with min abs_diff
    df_eeg_closest = pd.concat([df_eeg_closest, min_abs_diff])

# create a list of all folders in the given directory
pptx_list = os.listdir(Path(pptx_dir))
pptx_list2 = []
# run over pptx_list
for pptxs in pptx_list:
    if pptxs == "DB_5869_13.06.12.pptx":
        continue
    # get df_eeg_closest row with matching code
    row = df_eeg_closest.loc[
        df_eeg_closest["code2"].str.contains(pptxs.split("_")[1])
    ].iloc[0]
    # if row['EEG_date'] is in pptxs
    if row["EEG_date"] in pptxs:
        # add pptxs to pptx_list2
        pptx_list2.append(pptxs)

left, top, width, height = 0, 0, 1000, 500
# iterate over each folder in the list
for pptxs in pptx_list2:
    pswe_pptx = Presentation(pptx_dir + "/" + pptxs)
    # get image from pswe_pptx slide 3
    slide = pswe_pptx.slides[2]
    # get image from slide
    pswe_image = slide.shapes[0].image
    filename = "pswe_image.png"
    # save to .png file
    with open(filename, "wb") as f:
        f.write(pswe_image.blob)
    code = pptxs.split("_")[1]
    # find the matching row in the excel file
    row = Epilepsy_patients.loc[Epilepsy_patients["code"].str.contains(code)]
    int_match = 0
    # run over slides
    for i, slide in enumerate(pptx_file.slides):
        try:
            # find from pptx_file the 3rd slide with the matching code
            if slide.shapes.title.text.split("_")[1] == code:
                int_match += 1
                if int_match == 2:
                    int_match = 0
                    # add pswe image to slide
                    slide.shapes.add_picture(filename, left, top)  # , width, height)
                    pswe_image = slide.shapes[2].image
                    filename2 = "new.png"
                    # save to .png file
                    with open(filename2, "wb") as f:
                        f.write(pswe_image.blob)
        except:
            continue
# save pptx_file
pptx_file.save("Epilepsy_individulas3.pptx")
